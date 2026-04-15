import os
import shutil
from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt

from ga4_client import _normalize_top_page_key


# 유입경로 채널 고정 순서 (슬라이드5 표 행 1~9에 대응)
CHANNEL_ORDER = [
    "Referral",
    "Organic Search",
    "Direct",
    "Organic Social",
    "Unassigned",
    "Paid Search",
    "Display",
    "Organic Video",
    "Cross-network",
]

HIGHLIGHT_FILL = RGBColor(255, 242, 218)
WHITE_FILL = RGBColor(255, 255, 255)
LANG_LABELS = {"ko": "국문", "en": "영문", "cn": "중문"}
PAGE_SEGMENT_LABELS = {
    "": "메인",
    "index": "메인",
    "main": "메인",
    "home": "메인",
    "homepage": "메인",
    "daesang wellife": "메인",
    "daesang welllife": "메인",
    "daesang": "메인",
    "대상웰라이프": "메인",
    "首页": "메인",
    "主页": "메인",
    "about": "회사소개",
    "about-us": "회사소개",
    "about us": "회사소개",
    "company": "회사개요",
    "company-overview": "회사개요",
    "company overview": "회사개요",
    "公司概要": "회사개요",
    "overview": "회사개요",
    "intro": "회사소개",
    "introduction": "회사소개",
    "company-introduction": "회사소개",
    "company introduction": "회사소개",
    "company profile": "회사소개",
    "회사 소개": "회사소개",
    "公司介绍": "회사소개",
    "会社概要": "회사소개",
    "news": "뉴스",
    "news-list": "뉴스",
    "news-view": "뉴스",
    "list": "뉴스",
    "view": "뉴스",
    "media": "미디어",
    "brand": "브랜드",
    "newcare": "뉴케어",
    "new-care": "뉴케어",
    "new care": "뉴케어",
    "뉴 케어": "뉴케어",
    "纽凯儿": "뉴케어",
    "location": "사업장 위치",
    "locations": "사업장 위치",
    "business-location": "사업장 위치",
    "business-locations": "사업장 위치",
    "business location": "사업장 위치",
    "business locations": "사업장 위치",
    "사업장위치": "사업장 위치",
    "营业场所位置": "사업장 위치",
    "ir": "IR",
    "r&d": "R&D",
    "notice": "공고",
    "announcement": "공고",
    "rd": "R&D",
    "r-d": "R&D",
    "rnd": "R&D",
    "research": "R&D",
    "center": "센터",
    "rd-center": "R&D센터 소개",
    "r-d-center": "R&D센터 소개",
    "r&d center": "R&D센터 소개",
    "r&d center introduction": "R&D센터 소개",
    "r&d센터": "R&D센터 소개",
    "r&d센터 소개": "R&D센터 소개",
    "support": "고객지원",
    "faq": "FAQ",
    "welllife-solution": "Welllife Solution",
    "balanced-nutrition": "균형영양식",
    "balanced nutrition": "균형영양식",
    "balance": "균형영양식",
    "protein": "올프로틴",
    "all-protein": "올프로틴",
    "all protein": "올프로틴",
    "health-functional-food": "건강기능식품",
    "healthy-functional-food": "건강기능식품",
    "health functional food": "건강기능식품",
    "functional food": "건강기능식품",
    "affiliate": "계열사 소개",
    "affiliate-company": "계열사 소개",
    "affiliates": "계열사 소개",
    "affiliate company": "계열사 소개",
    "affiliated company": "계열사 소개",
    "affiliated companies": "계열사 소개",
    "value": "가치체계",
    "values": "가치체계",
    "value system": "가치체계",
}


# ── 내부 유틸 ────────────────────────────────────────────────

def _get_shape(slide, name: str):
    """슬라이드에서 이름으로 도형을 찾는다."""
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    raise ValueError(f"도형 '{name}'을 슬라이드에서 찾을 수 없습니다.")


def _set_cell(table, row: int, col: int, value) -> None:
    """
    표 셀의 텍스트를 변경한다.
    기존 런(run)의 서식(폰트, 색상 등)을 유지하면서 텍스트만 교체한다.
    """
    cell = table.cell(row, col)
    tf = cell.text_frame

    # 단락이 여러 개면 첫 번째만 사용
    while len(tf.paragraphs) > 1:
        p = tf.paragraphs[-1]._p
        p.getparent().remove(p)

    display_value = f"{value:,}" if isinstance(value, int) else str(value)

    para = tf.paragraphs[0]
    if para.runs:
        # 첫 번째 런에 값 입력, 나머지 런은 비움
        para.runs[0].text = display_value
        for run in para.runs[1:]:
            run.text = ""
    else:
        para.text = display_value


def _set_textbox(slide, shape_name: str, text: str) -> None:
    """텍스트박스 내용을 교체한다. 기존 문단을 재사용하되 남는 문단은 제거한다."""
    shape = _get_shape(slide, shape_name)
    tf = shape.text_frame
    lines = text.splitlines() or [""]

    while len(tf.paragraphs) < len(lines):
        source_para = tf.paragraphs[-1]
        source_run = source_para.runs[0] if source_para.runs else None
        new_para = tf.add_paragraph()

        if source_para._p.pPr is not None:
            new_para._p.insert(0, deepcopy(source_para._p.pPr))

        if source_run is not None:
            new_run = new_para.add_run()
            if source_run._r.rPr is not None:
                new_run._r.insert(0, deepcopy(source_run._r.rPr))

    while len(tf.paragraphs) > len(lines):
        p = tf.paragraphs[-1]._p
        p.getparent().remove(p)

    for para, line in zip(tf.paragraphs, lines):
        if not para.runs:
            para.add_run()

        first_run = para.runs[0]
        first_run.text = line

        # Keep the first run so its font/size/bold/color survive, but remove
        # stale runs and manual line-break nodes left from older reports.
        first_run_element = first_run._r
        for child in list(para._p):
            tag_name = child.tag.rsplit("}", 1)[-1]
            if tag_name == "pPr":
                continue
            if child is first_run_element:
                continue
            para._p.remove(child)


def _set_cell_bold(table, row: int, col: int, bold: bool) -> None:
    """표 셀 텍스트의 bold 여부를 설정한다."""
    cell = table.cell(row, col)
    for para in cell.text_frame.paragraphs:
        if not para.runs:
            para.text = para.text
        for run in para.runs:
            run.font.bold = bold


def _set_row_bold(table, row: int, bold: bool = True) -> None:
    """표의 특정 행 전체 bold 여부를 설정한다."""
    for col in range(len(table.columns)):
        _set_cell_bold(table, row, col, bold)


def _set_cell_font(table, row: int, col: int, font_name: str, font_size_pt: float) -> None:
    """표 셀 텍스트의 폰트명/크기를 강제 적용한다."""
    cell = table.cell(row, col)
    tf = cell.text_frame
    for para in tf.paragraphs:
        if not para.runs:
            para.add_run()
        for run in para.runs:
            run.font.name = font_name
            run.font.size = Pt(font_size_pt)


def _rank_eq(values: list[int]) -> list[int]:
    """Excel RANK.EQ(내림차순)과 동일한 순위를 반환한다."""
    return [1 + sum(1 for other in values if other > value) for value in values]


def _set_cell_fill(table, row: int, col: int, color: RGBColor) -> None:
    """표 셀 배경색을 설정한다."""
    cell = table.cell(row, col)
    cell.fill.solid()
    cell.fill.fore_color.rgb = color


def _set_cell_emphasis(table, row: int, col: int, bold: bool, color: RGBColor) -> None:
    """표 셀의 강조 서식을 한 번에 설정한다."""
    _set_cell_bold(table, row, col, bold)
    _set_cell_fill(table, row, col, color)


def _get_cell_int(table, row: int, col: int) -> int:
    """표 셀에서 정수 값을 읽는다. 숫자가 아닌 경우 0 반환."""
    cell = table.cell(row, col)
    tf = cell.text_frame
    text = tf.paragraphs[0].text.strip().replace(",", "")
    try:
        return int(text)
    except (ValueError, AttributeError):
        return 0


def _format_rate(rate: float) -> str:
    """증가율을 '▲X%' / '▼X%' 형식으로 포맷한다."""
    symbol = "▲" if rate >= 0 else "▼"
    return f"{symbol}{abs(rate):.1f}%"


def _format_period(start_date: str, end_date: str) -> str:
    """YYYY-MM-DD 기간을 '기간 : M/D~M/D' 형식으로 변환한다."""
    import datetime

    start = datetime.datetime.strptime(start_date, "%Y-%m-%d")
    end = datetime.datetime.strptime(end_date, "%Y-%m-%d")
    return f"기간 : {start.month}/{start.day}~{end.month}/{end.day}"


def _format_duration(seconds: float | int | None) -> str:
    """초 단위 체류 시간을 00:00:00 형식으로 변환한다."""
    # 소수점은 반올림하지 않고 버림 처리
    total_seconds = int(seconds or 0)
    hours = total_seconds // 3600
    minutes = (total_seconds % 3600) // 60
    secs = total_seconds % 60
    return f"{hours:02d}:{minutes:02d}:{secs:02d}"


def _normalize_page_segment(segment: str) -> str:
    """URL segment를 사람이 읽기 쉬운 페이지명으로 변환한다."""
    from urllib.parse import unquote

    cleaned = unquote(segment).strip().strip("/")
    if not cleaned:
        return "메인"

    lowered = cleaned.lower().replace("_", "-")
    if lowered in PAGE_SEGMENT_LABELS:
        return PAGE_SEGMENT_LABELS[lowered]

    words = [PAGE_SEGMENT_LABELS.get(part, part) for part in lowered.split("-")]
    return " ".join(words)


def _translate_page_path(path: str) -> str:
    """GA4 pagePath를 보고서용 페이지 경로로 변환한다."""
    return _normalize_top_page_key(path)


def _summarize_page_title(path: str) -> str:
    """상단 TOP 5 요약에 사용할 짧은 페이지명을 반환한다."""
    raw = (path or "").strip()
    if raw == "메인":
        return raw
    return raw.split(" > ")[-1] if raw else "-"


def _format_top_pages_summary(data: dict) -> str:
    """6p 상단 인기 페이지 TOP 5와 평균 체류 시간 문구를 생성한다."""
    lines = []
    for lang in ["ko", "en", "cn"]:
        pages = data[lang].get("top_pages", [])[:5]
        titles = [_summarize_page_title(item.get("page", "")) for item in pages]
        summary = " > ".join(titles) if titles else "-"
        lines.append(f"- {LANG_LABELS[lang]}: {summary} 순")

    durations = ", ".join(
        f"{LANG_LABELS[lang]} {_format_duration(data[lang].get('avg_engagement'))}"
        for lang in ["ko", "en", "cn"]
    )
    lines.append(f"- 평균 체류 시간 : {durations}")
    return "\n".join(lines)


def _fill_detail_period(prs: Presentation, start_date: str, end_date: str) -> None:
    """상세페이지 레이아웃 우측 상단 기간 텍스트를 업데이트한다."""
    period_text = _format_period(start_date, end_date)
    for layout in prs.slide_layouts:
        if layout.name == "상세페이지":
            _set_textbox(layout, "TextBox 1", period_text)
            return


# ── 슬라이드별 입력 함수 ─────────────────────────────────────

def _fill_cover(slide, client_name: str, year: int, month: int) -> None:
    """슬라이드1: 표지 제목과 날짜를 업데이트한다."""
    import datetime
    today = datetime.date.today()
    title_text = f"{client_name} 웹사이트\n{month}월 운영보고서"
    date_text = f"더위버크리에이티브\n{today.year}.{today.month:02d}.{today.day:02d}"
    _set_textbox(slide, "TextBox 2", title_text)
    _set_textbox(slide, "TextBox 3", date_text)


def _fill_metric_table(slide, month: int, ko_val: int, en_val: int, cn_val: int) -> dict:
    """
    슬라이드3(사용자수) / 슬라이드4(페이지뷰수):
    보고월 행에 국/영/중 수치, 소계(col4), 증가율(col6)을 입력하고
    합계 행(행14)도 국/영/중/소계 합산하여 업데이트한다.
    col5(2025실적)는 건드리지 않는다.

    Returns:
        {"ko": float|None, "en": float|None, "cn": float|None}
        전월 데이터가 없는 경우(1월) None 반환
    """
    table = _get_shape(slide, "표 6").table
    row_idx = month + 1  # 1월=2, 2월=3, 3월=4 ...

    # 현재 월 수치 입력
    _set_cell(table, row_idx, 1, ko_val)
    _set_cell(table, row_idx, 2, en_val)
    _set_cell(table, row_idx, 3, cn_val)

    # 소계(col4) = 국문 + 영문 + 중문
    subtotal = ko_val + en_val + cn_val
    _set_cell(table, row_idx, 4, subtotal)

    # 2025 실적(col5)이 0이면 증가율 계산을 수행하지 않는다.
    base_2025 = _get_cell_int(table, row_idx, 5)

    # 전월 대비 증가율 계산 (1월=row2는 전월 없음)
    growth_rates = {"ko": None, "en": None, "cn": None}
    prev_row_idx = row_idx - 1
    if prev_row_idx >= 2 and base_2025 != 0:  # 2월 이후 + 2025 실적 존재 시 계산
        prev_ko = _get_cell_int(table, prev_row_idx, 1)
        prev_en = _get_cell_int(table, prev_row_idx, 2)
        prev_cn = _get_cell_int(table, prev_row_idx, 3)
        prev_subtotal = prev_ko + prev_en + prev_cn

        def _calc(curr, prev):
            return (curr - prev) / prev * 100 if prev != 0 else None

        growth_rates = {
            "ko": _calc(ko_val, prev_ko),
            "en": _calc(en_val, prev_en),
            "cn": _calc(cn_val, prev_cn),
        }

        # 소계 기준 증가율을 col6에 입력
        total_rate = _calc(subtotal, prev_subtotal)
        if total_rate is not None:
            _set_cell(table, row_idx, 6, _format_rate(total_rate))
    else:
        _set_cell(table, row_idx, 6, "-")

    # 합계 행(행14): 행2~13(1월~12월) 합산
    ko_total = en_total = cn_total = 0
    for r in range(2, 14):
        ko_total += _get_cell_int(table, r, 1)
        en_total += _get_cell_int(table, r, 2)
        cn_total += _get_cell_int(table, r, 3)
    _set_cell(table, 14, 1, ko_total)
    _set_cell(table, 14, 2, en_total)
    _set_cell(table, 14, 3, cn_total)
    _set_cell(table, 14, 4, ko_total + en_total + cn_total)
    _set_row_bold(table, 14, True)
    _position_month_highlight(slide, table, row_idx)

    return growth_rates


def _position_month_highlight(slide, table, row_idx: int) -> None:
    """3p/4p의 빨간 박스를 해당 월 행으로 이동시킨다."""
    highlight = None
    for shape in slide.shapes:
        if shape.name == "직사각형 5":
            highlight = shape
            break
    if highlight is None:
        return

    # 템플릿 표의 실제 행 위치를 기준으로 빨간 박스를 수직 정렬
    row_top = table._graphic_frame.top + sum(table.rows[r].height for r in range(row_idx))
    row_height = table.rows[row_idx].height
    highlight.top = row_top + max(0, (row_height - highlight.height) // 2)


def _fill_summary_textbox(slide, metric_name: str, growth_rates: dict) -> None:
    """
    슬라이드3~4: TextBox 1(요약문)에 언어별 전월 대비 증감률 문구를 입력한다.
    단락이 3개인 템플릿 기준으로 국/영/중 순서로 업데이트한다.
    growth_rates: {"ko": float|None, "en": float|None, "cn": float|None}
    """
    labels = {"ko": "국문", "en": "영문", "cn": "중문"}
    lines = []
    for lang in ["ko", "en", "cn"]:
        rate = growth_rates[lang]
        if rate is None:
            lines.append(f"- {labels[lang]}: 전월 기준 {metric_name} 정보 없음")
        else:
            direction = "증가" if rate >= 0 else "감소"
            lines.append(f"- {labels[lang]}: 전월 기준 {metric_name} {abs(rate):.0f}% {direction}")

    shape = _get_shape(slide, "TextBox 1")
    _set_textbox(slide, shape.name, "\n".join(lines))


def _fill_channel_table(
    slide,
    ko_channels: list[dict],
    en_channels: list[dict],
    cn_channels: list[dict],
) -> None:
    """
    슬라이드5: 유입경로 표에 채널별 세션수를 입력한다.
    비율 열(col2/4/6)은 PPT에 수식이 없으므로 직접 계산하여 업데이트한다.
    합계 행(행10)도 함께 업데이트한다.
    """
    table = _get_shape(slide, "표 3").table

    def build_map(channels):
        return {item["channel"]: item["sessions"] for item in channels}

    ko_map = build_map(ko_channels)
    en_map = build_map(en_channels)
    cn_map = build_map(cn_channels)

    ko_total = sum(ko_map.get(ch, 0) for ch in CHANNEL_ORDER)
    en_total = sum(en_map.get(ch, 0) for ch in CHANNEL_ORDER)
    cn_total = sum(cn_map.get(ch, 0) for ch in CHANNEL_ORDER)
    lang_specs = [
        (ko_map, ko_total, 1, 2),
        (en_map, en_total, 3, 4),
        (cn_map, cn_total, 5, 6),
    ]
    max_values = {
        value_col: max((channel_map.get(ch, 0) for ch in CHANNEL_ORDER), default=0)
        for channel_map, _total, value_col, _rate_col in lang_specs
    }

    for i, channel in enumerate(CHANNEL_ORDER):
        row_idx = i + 1
        ko_val = ko_map.get(channel, 0)
        en_val = en_map.get(channel, 0)
        cn_val = cn_map.get(channel, 0)

        # 세션수 입력
        _set_cell(table, row_idx, 1, ko_val)
        _set_cell(table, row_idx, 3, en_val)
        _set_cell(table, row_idx, 5, cn_val)

        # 비율 계산 및 입력 (PPT는 수식 없음)
        _set_cell(table, row_idx, 2, f"{ko_val/ko_total*100:.2f}%" if ko_total else "0.00%")
        _set_cell(table, row_idx, 4, f"{en_val/en_total*100:.2f}%" if en_total else "0.00%")
        _set_cell(table, row_idx, 6, f"{cn_val/cn_total*100:.2f}%" if cn_total else "0.00%")

        for channel_map, _total, value_col, rate_col in lang_specs:
            value = channel_map.get(channel, 0)
            is_max = max_values[value_col] > 0 and value == max_values[value_col]
            fill_color = HIGHLIGHT_FILL if is_max else WHITE_FILL
            _set_cell_emphasis(table, row_idx, value_col, is_max, fill_color)
            _set_cell_emphasis(table, row_idx, rate_col, is_max, fill_color)

    # 합계 행 (행 10)
    _set_cell(table, 10, 1, ko_total)
    _set_cell(table, 10, 3, en_total)
    _set_cell(table, 10, 5, cn_total)
    _set_cell(table, 10, 2, "100.00%")
    _set_cell(table, 10, 4, "100.00%")
    _set_cell(table, 10, 6, "100.00%")
    _set_row_bold(table, 10, True)


def _fill_top_pages_table(slide, table_name: str, pages: list[dict]) -> None:
    """
    슬라이드6: 인기 페이지 표(국/영/중)에 TOP 10 데이터를 입력한다.
    행0(언어 헤더), 행1(컬럼 헤더)는 건드리지 않는다.
    """
    table = _get_shape(slide, table_name).table
    pageviews_values = [int(item["pageviews"]) for item in pages[:10]]
    ranks = _rank_eq(pageviews_values) if pageviews_values else []

    for i in range(10):
        row_idx = i + 2  # 데이터는 2행부터
        if i < len(pages):
            _set_cell(table, row_idx, 0, ranks[i] if i < len(ranks) else "-")
            _set_cell(table, row_idx, 1, pages[i]["page"])
            if table_name == "표 3":
                _set_cell_font(table, row_idx, 1, "맑은 고딕", 10)
            _set_cell(table, row_idx, 2, pages[i]["pageviews"])
        else:
            _set_cell(table, row_idx, 0, "-")
            _set_cell(table, row_idx, 1, "-")
            if table_name == "표 3":
                _set_cell_font(table, row_idx, 1, "맑은 고딕", 10)
            _set_cell(table, row_idx, 2, "-")


# ── 공개 인터페이스 ──────────────────────────────────────────

def write_report(
    template_path: str,
    output_dir: str,
    client_name: str,
    year: int,
    month: int,
    data: dict,
    start_date: str | None = None,
    end_date: str | None = None,
) -> str:
    """
    GA4 데이터를 PPT 템플릿에 입력하고 output 폴더에 저장한다.

    Args:
        template_path: 템플릿 pptx 파일 경로
        output_dir:    결과 파일 저장 폴더
        client_name:   고객사명 (파일명 및 표지에 사용)
        year:          보고 연도 (예: 2026)
        month:         보고월 (예: 3)
        data: {
            "ko": {"users": int, "sessions": int, "pageviews": int,
                   "channels": [...], "top_pages": [...]},
            "en": {...},
            "cn": {...},
        }

    Returns:
        저장된 파일 경로
    """
    output_path = os.path.join(output_dir, f"{client_name}_{year}{month:02d}_report.pptx")
    Path(output_dir).mkdir(parents=True, exist_ok=True)

    shutil.copy2(template_path, output_path)
    prs = Presentation(output_path)

    slides = prs.slides  # 0-based index
    if start_date and end_date:
        _fill_detail_period(prs, start_date, end_date)

    # 슬라이드1: 표지
    _fill_cover(slides[0], client_name, year, month)

    # 슬라이드3: 사용자수
    user_rates = _fill_metric_table(
        slides[2], month,
        data["ko"]["users"],
        data["en"]["users"],
        data["cn"]["users"],
    )
    _fill_summary_textbox(slides[2], "사용자수", user_rates)

    # 슬라이드4: 페이지뷰수
    pv_rates = _fill_metric_table(
        slides[3], month,
        data["ko"]["pageviews"],
        data["en"]["pageviews"],
        data["cn"]["pageviews"],
    )
    _fill_summary_textbox(slides[3], "페이지뷰수", pv_rates)

    # 슬라이드5: 유입경로
    _fill_channel_table(
        slides[4],
        data["ko"]["channels"],
        data["en"]["channels"],
        data["cn"]["channels"],
    )

    # 슬라이드6: 인기 페이지
    _set_textbox(slides[5], "TextBox 14", _format_top_pages_summary(data))
    _fill_top_pages_table(slides[5], "표 3", data["ko"]["top_pages"])
    _fill_top_pages_table(slides[5], "표 4", data["en"]["top_pages"])
    _fill_top_pages_table(slides[5], "표 7", data["cn"]["top_pages"])

    prs.save(output_path)
    print(f"저장 완료: {output_path}")
    return output_path
