import argparse
import calendar
import json
from datetime import datetime
from pathlib import Path
from typing import Any

from dotenv import load_dotenv
from openpyxl import load_workbook
from pptx import Presentation
from pptx.util import Pt
from google.analytics.data_v1beta.types import (
    DateRange,
    Dimension,
    Metric,
    OrderBy,
    RunReportRequest,
)

import ppt_gen
from ga4_client import (
    fetch_avg_engagement,
    fetch_channel_sessions,
    fetch_summary,
    fetch_top_pages,
    get_ga4_client,
)

load_dotenv()


BASE_DIR = Path(__file__).resolve().parent
CONFIG_PATH = BASE_DIR / "config" / "clients.json"
TEMPLATES_DIR = BASE_DIR / "templates"
OUTPUT_DIR = BASE_DIR / "output"
BASELINE_DIR = BASE_DIR / "config" / "annual_baseline"
LOG_DIR = BASE_DIR / "logs"
LOG_PATH = LOG_DIR / "run_log.txt"
LANGUAGES = ("ko", "en", "cn")
IA_FILENAME = "대상웰라이프_IA.xlsx"
IA_LANG_COLS = {"ko": 5, "en": 6, "cn": 7}
MAIN_EQUIVALENT_ALIASES = {
    "대상웰라이프",
    "daesangwellife",
    "大象wellife",
    "main",
    "home",
    "homepage",
    "index",
    "kr",
    "en",
    "cn",
    "메인",
}


def parse_args() -> argparse.Namespace:
    """CLI 인자를 정의하고 파싱한다."""
    parser = argparse.ArgumentParser(
        description="GA4 데이터를 조회해 엑셀/PPT 운영 보고서를 생성합니다."
    )
    parser.add_argument("--client", required=True, help="clients.json에 등록된 고객사명")
    parser.add_argument("--start", required=True, help="조회 시작일 (YYYYMMDD)")
    parser.add_argument("--end", required=True, help="조회 종료일 (YYYYMMDD)")
    parser.add_argument(
        "--report-month",
        help="보고서 기준 월 (YYYYMM). 미입력 시 --end의 연월을 사용",
    )
    return parser.parse_args()


def parse_date(date_str: str) -> datetime:
    """YYYYMMDD 문자열을 datetime으로 변환한다."""
    try:
        return datetime.strptime(date_str, "%Y%m%d")
    except ValueError as exc:
        raise ValueError(f"잘못된 날짜 형식입니다: {date_str} (예: 20260301)") from exc


def format_ga4_date(date_str: str) -> str:
    """YYYYMMDD 문자열을 GA4용 YYYY-MM-DD 형식으로 변환한다."""
    return parse_date(date_str).strftime("%Y-%m-%d")


def format_report_month(date_value: datetime) -> str:
    """보고서 파일명에 사용할 YYYYMM 문자열을 반환한다."""
    return date_value.strftime("%Y%m")


def parse_report_month(report_month_str: str) -> datetime:
    """YYYYMM 문자열을 보고서 기준 월의 1일 datetime으로 변환한다."""
    try:
        return datetime.strptime(report_month_str, "%Y%m")
    except ValueError as exc:
        raise ValueError(
            f"잘못된 보고서월 형식입니다: {report_month_str} (예: 202603)"
        ) from exc


def get_previous_month(date_value: datetime) -> datetime:
    """주어진 날짜 기준 전월 1일을 반환한다."""
    if date_value.month == 1:
        return date_value.replace(year=date_value.year - 1, month=12, day=1)
    return date_value.replace(month=date_value.month - 1, day=1)


def shift_month(date_value: datetime, delta: int) -> datetime:
    """월 단위로 이동하되 day는 가능한 범위 내에서 유지한다."""
    year = date_value.year + ((date_value.month - 1 + delta) // 12)
    month = (date_value.month - 1 + delta) % 12 + 1
    max_day = calendar.monthrange(year, month)[1]
    day = min(date_value.day, max_day)
    return date_value.replace(year=year, month=month, day=day)


def load_client_config(client_name: str) -> dict:
    """clients.json에서 고객사 설정을 읽어 반환한다."""
    if not CONFIG_PATH.exists():
        raise FileNotFoundError(f"고객사 설정 파일을 찾을 수 없습니다: {CONFIG_PATH}")

    with CONFIG_PATH.open("r", encoding="utf-8") as file:
        config = json.load(file)

    clients = config.get("clients", [])
    for client in clients:
        if client.get("name") == client_name:
            return client

    available_names = [client.get("name", "") for client in clients if client.get("name")]
    available_text = ", ".join(available_names) if available_names else "(등록된 고객사 없음)"
    raise ValueError(
        f"clients.json에서 고객사 '{client_name}' 설정을 찾을 수 없습니다. "
        f"등록된 고객사: {available_text}"
    )


def resolve_template_path(filename: str) -> str:
    """템플릿 파일명을 절대경로로 변환하고 존재 여부를 확인한다."""
    template_path = TEMPLATES_DIR / filename
    if not template_path.exists():
        raise FileNotFoundError(f"템플릿 파일을 찾을 수 없습니다: {template_path}")
    return str(template_path)


def build_report_filename(client_name: str, report_month: str, extension: str) -> str:
    """월별 산출물 파일명을 만든다."""
    return f"{client_name}_{report_month}_report.{extension}"


def get_client_output_dir(client_name: str) -> Path:
    """고객사별 output 하위 폴더 경로를 반환한다."""
    return OUTPUT_DIR / client_name


def resolve_output_path(client_name: str, report_month: str, extension: str) -> Path:
    """고객사별 월 산출물 절대경로를 반환한다."""
    return get_client_output_dir(client_name) / build_report_filename(
        client_name, report_month, extension
    )


def _normalize_monthly_totals(values: Any, metric_name: str) -> list[int]:
    """월별(1~12월) 합계 리스트를 int[12]로 정규화한다."""
    if not isinstance(values, list) or len(values) != 12:
        raise ValueError(f"annual baseline '{metric_name}' must be a list with 12 items.")
    try:
        return [int(value) for value in values]
    except (TypeError, ValueError) as exc:
        raise ValueError(f"annual baseline '{metric_name}' contains a non-integer value.") from exc


def load_annual_baseline(client_name: str, report_month_dt: datetime) -> dict[str, Any]:
    """
    전년도(보고월 기준 year-1) 월별 합계 baseline을 로드한다.

    파일 형식:
    {
      "users_total_monthly": [12개 정수],
      "pageviews_total_monthly": [12개 정수]
    }
    """
    baseline_year = report_month_dt.year - 1
    baseline_path = BASELINE_DIR / client_name / f"{baseline_year}.json"
    if not baseline_path.exists():
        raise FileNotFoundError(
            f"annual baseline file not found: client={client_name} "
            f"baseline_year={baseline_year} path={baseline_path}"
        )

    with baseline_path.open("r", encoding="utf-8") as file:
        raw = json.load(file)

    users = _normalize_monthly_totals(raw.get("users_total_monthly"), "users_total_monthly")
    pageviews = _normalize_monthly_totals(
        raw.get("pageviews_total_monthly"), "pageviews_total_monthly"
    )
    return {
        "year": baseline_year,
        "path": str(baseline_path),
        "users_total_monthly": users,
        "pageviews_total_monthly": pageviews,
    }


def _set_ppt_table_cell_int(
    table,
    row: int,
    col: int,
    value: int | str | None,
    *,
    force_bold: bool = False,
    font_name: str | None = None,
    font_size_pt: int | None = None,
) -> None:
    """PPT 테이블 셀 텍스트를 정수 값으로 설정한다(기존 폰트/스타일 최대 유지)."""
    cell = table.cell(row, col)
    tf = cell.text_frame
    if value is None:
        display = "-"
    elif isinstance(value, str):
        display = value
    else:
        display = f"{int(value):,}"

    def _apply_run_style(run) -> None:
        if force_bold:
            run.font.bold = True
        if font_name:
            run.font.name = font_name
        if font_size_pt:
            run.font.size = Pt(font_size_pt)
    if tf.paragraphs:
        para = tf.paragraphs[0]
        if para.runs:
            para.runs[0].text = display
            _apply_run_style(para.runs[0])
            for run in para.runs[1:]:
                run.text = ""
                _apply_run_style(run)
        else:
            para.text = display
            if para.runs:
                _apply_run_style(para.runs[0])
    else:
        tf.text = display
        if tf.paragraphs and tf.paragraphs[0].runs:
            _apply_run_style(tf.paragraphs[0].runs[0])


def _find_shape_by_name(slide, shape_name: str):
    """슬라이드에서 지정된 이름의 도형을 찾는다."""
    for shape in slide.shapes:
        if shape.name == shape_name:
            return shape
    raise ValueError(f"PPT shape '{shape_name}' not found on slide.")


def _calc_growth_rate(curr: int, prev: int) -> float | None:
    """전월 대비 증감율을 계산한다."""
    if prev == 0:
        return None
    return (curr - prev) / prev * 100


def _format_growth_rate(rate: float) -> str:
    """증감율 텍스트를 +x.x% / -x.x% 형식으로 반환한다."""
    sign = "+" if rate >= 0 else "-"
    return f"{sign}{abs(rate):.1f}%"


def _set_textbox_multiline(slide, shape_name: str, lines: list[str]) -> None:
    """TextBox 내용을 여러 줄로 교체한다."""
    shape = _find_shape_by_name(slide, shape_name)
    tf = shape.text_frame
    text_lines = lines if lines else [""]

    while len(tf.paragraphs) < len(text_lines):
        tf.add_paragraph()
    while len(tf.paragraphs) > len(text_lines):
        p = tf.paragraphs[-1]._p
        p.getparent().remove(p)

    for para, line in zip(tf.paragraphs, text_lines):
        para.text = line
        for run in para.runs:
            run.font.name = "맑은 고딕"
            run.font.size = Pt(14)


def apply_growth_overrides_to_generated_ppt(
    ppt_path: str,
    report_month_dt: datetime,
    monthly_summary_series: dict[int, dict[str, dict[str, int | None]]],
    annual_baseline: dict[str, Any],
) -> None:
    """
    생성된 PPT의 3p/4p 증감율(표 col6 + 요약 TextBox 1)을
    월별 GA4 데이터 기준으로 재계산해 덮어쓴다.
    """
    report_month = report_month_dt.month
    if report_month < 2:
        return

    curr = monthly_summary_series.get(report_month)
    prev = monthly_summary_series.get(report_month - 1)
    if not curr or not prev:
        return

    prs = Presentation(ppt_path)
    users_slide = prs.slides[2]
    pv_slide = prs.slides[3]
    row_idx = report_month + 1

    def _apply_metric(slide, key: str, metric_name: str, baseline_values: list[int]) -> None:
        table = _find_shape_by_name(slide, "표 6").table
        curr_values = [curr["ko"].get(key), curr["en"].get(key), curr["cn"].get(key)]
        curr_total = (
            None
            if any(value is None for value in curr_values)
            else int(curr_values[0]) + int(curr_values[1]) + int(curr_values[2])
        )
        baseline_month_value = int(baseline_values[report_month - 1])
        total_rate = (
            None
            if curr_total is None
            else _calc_growth_rate(int(curr_total), baseline_month_value)
        )
        if total_rate is None:
            _set_ppt_table_cell_int(
                table,
                row_idx,
                6,
                "-",
                font_name="맑은 고딕",
                font_size_pt=14,
            )
        else:
            _set_ppt_table_cell_int(
                table,
                row_idx,
                6,
                _format_growth_rate(total_rate),
                font_name="맑은 고딕",
                font_size_pt=14,
            )

        lines = []
        for lang, label in [("ko", "국문"), ("en", "영문"), ("cn", "중문")]:
            curr_lang_value = curr[lang].get(key)
            prev_lang_value = prev[lang].get(key)
            if curr_lang_value is None or prev_lang_value is None:
                lang_rate = None
            else:
                lang_rate = _calc_growth_rate(int(curr_lang_value), int(prev_lang_value))
            if lang_rate is None:
                lines.append(f"- {label}: -")
            else:
                direction = "증가" if lang_rate >= 0 else "감소"
                lines.append(f"- {label}: 전월 기준 {metric_name} {abs(lang_rate):.0f}% {direction}")
        _set_textbox_multiline(slide, "TextBox 1", lines)

    _apply_metric(users_slide, "users", "사용자수", annual_baseline["users_total_monthly"])
    _apply_metric(
        pv_slide,
        "pageviews",
        "페이지뷰수",
        annual_baseline["pageviews_total_monthly"],
    )
    prs.save(ppt_path)


def apply_annual_baseline_to_ppt_base(
    ppt_base_path: str,
    client_output_dir: Path,
    report_month: str,
    annual_baseline: dict[str, Any],
    monthly_summary_series: dict[int, dict[str, dict[str, int | None]]],
) -> str:
    """
    3p/4p 표의 전년도 실적(열 index 5)에 baseline(월별 합계) 값을 채운
    임시 base PPT를 생성하고 경로를 반환한다.
    """
    prs = Presentation(ppt_base_path)
    if len(prs.slides) < 4:
        raise ValueError("PPT template has fewer slides than expected for baseline injection.")

    users_values = annual_baseline["users_total_monthly"]
    pageviews_values = annual_baseline["pageviews_total_monthly"]

    users_table = _find_shape_by_name(prs.slides[2], "표 6").table
    pageviews_table = _find_shape_by_name(prs.slides[3], "표 6").table

    for month_idx, value in enumerate(users_values, start=1):
        _set_ppt_table_cell_int(users_table, month_idx + 1, 5, value, force_bold=True)
    _set_ppt_table_cell_int(users_table, 14, 5, sum(int(v) for v in users_values), force_bold=True)

    for month_idx, value in enumerate(pageviews_values, start=1):
        _set_ppt_table_cell_int(pageviews_table, month_idx + 1, 5, value, force_bold=True)
    _set_ppt_table_cell_int(
        pageviews_table,
        14,
        5,
        sum(int(v) for v in pageviews_values),
        force_bold=True,
    )

    # 3/4p 월별 누적(1월~보고월) 국/영/중문 + 합계를 GA4 실측값으로 채운다.
    for month_idx in sorted(monthly_summary_series.keys()):
        row_idx = month_idx + 1
        month_data = monthly_summary_series[month_idx]

        users_ko = month_data["ko"]["users"]
        users_en = month_data["en"]["users"]
        users_cn = month_data["cn"]["users"]
        _set_ppt_table_cell_int(users_table, row_idx, 1, users_ko)
        _set_ppt_table_cell_int(users_table, row_idx, 2, users_en)
        _set_ppt_table_cell_int(users_table, row_idx, 3, users_cn)
        users_total = (
            None
            if users_ko is None or users_en is None or users_cn is None
            else int(users_ko) + int(users_en) + int(users_cn)
        )
        _set_ppt_table_cell_int(users_table, row_idx, 4, users_total)

        pv_ko = month_data["ko"]["pageviews"]
        pv_en = month_data["en"]["pageviews"]
        pv_cn = month_data["cn"]["pageviews"]
        _set_ppt_table_cell_int(pageviews_table, row_idx, 1, pv_ko)
        _set_ppt_table_cell_int(pageviews_table, row_idx, 2, pv_en)
        _set_ppt_table_cell_int(pageviews_table, row_idx, 3, pv_cn)
        pv_total = (
            None
            if pv_ko is None or pv_en is None or pv_cn is None
            else int(pv_ko) + int(pv_en) + int(pv_cn)
        )
        _set_ppt_table_cell_int(pageviews_table, row_idx, 4, pv_total)

    runtime_base_path = client_output_dir / f"_runtime_base_{report_month}.pptx"
    prs.save(str(runtime_base_path))
    return str(runtime_base_path)


def resolve_base_report_path(
    client_name: str,
    report_month_dt: datetime,
    template_filename: str,
    extension: str,
) -> tuple[str, str]:
    """
    입력 베이스 파일 경로와 출처를 반환한다.

    우선순위:
    1. output/{client}/ 전월 파일
    2. 기존 flat output/ 전월 파일 (하위 호환)
    3. templates/ 템플릿 파일
    """
    prev_month = format_report_month(get_previous_month(report_month_dt))
    client_prev_path = resolve_output_path(client_name, prev_month, extension)
    if client_prev_path.exists():
        return str(client_prev_path), "prev_month_file"

    legacy_prev_path = OUTPUT_DIR / build_report_filename(client_name, prev_month, extension)
    if legacy_prev_path.exists():
        return str(legacy_prev_path), "prev_month_file_legacy"

    return resolve_template_path(template_filename), "template"


def fetch_summary_safe(
    property_id: str,
    start_date: str,
    end_date: str,
    *,
    context: str,
) -> dict[str, int | bool]:
    """
    GA4 summary 조회를 안전하게 수행한다.
    응답 rows가 비어 IndexError가 발생하면 0 값으로 대체한다.
    """
    try:
        summary = fetch_summary(property_id, start_date, end_date)
        summary["_missing"] = False
        return summary
    except IndexError:
        write_log(
            "GA4_EMPTY_SUMMARY "
            f"context={context} property_id={property_id} start={start_date} end={end_date}"
        )
        return {"users": 0, "sessions": 0, "pageviews": 0, "_missing": True}


def collect_ga4_data(client_config: dict, start_date: str, end_date: str) -> dict:
    """국문/영문/중문 GA4 데이터를 조회해 보고서용 딕셔너리로 조립한다."""
    property_ids = client_config.get("ga4_property_ids", {})
    missing_languages = [lang for lang in LANGUAGES if not property_ids.get(lang)]
    if missing_languages:
        missing = ", ".join(missing_languages)
        raise ValueError(f"GA4 속성 ID가 비어 있습니다: {missing}")

    data = {}
    for lang in LANGUAGES:
        property_id = property_ids[lang]
        summary = fetch_summary_safe(
            property_id,
            start_date,
            end_date,
            context=f"current_month_{lang}",
        )
        channels = fetch_channel_sessions(property_id, start_date, end_date)
        top_pages = fetch_top_pages(property_id, start_date, end_date)
        avg_engagement = fetch_avg_engagement(property_id, start_date, end_date)

        data[lang] = {
            "users": int(summary["users"]),
            "sessions": int(summary["sessions"]),
            "pageviews": int(summary["pageviews"]),
            "channels": channels,
            "top_pages": top_pages,
            "avg_engagement": avg_engagement,
        }

    return data


def collect_monthly_summary_series(
    client_config: dict,
    start_dt: datetime,
    end_dt: datetime,
    report_month_dt: datetime,
) -> dict[int, dict[str, dict[str, int | None]]]:
    """
    1월~보고월까지 월별 GA4 summary(totalUsers/sessions/pageviews)를 수집한다.

    반환 형식:
    {
      1: {"ko": {"users":.., "sessions":.., "pageviews":..}, "en": {...}, "cn": {...}},
      ...
      report_month: {...}
    }
    """
    property_ids = client_config.get("ga4_property_ids", {})
    missing_languages = [lang for lang in LANGUAGES if not property_ids.get(lang)]
    if missing_languages:
        missing = ", ".join(missing_languages)
        raise ValueError(f"GA4 속성 ID가 비어 있습니다: {missing}")

    report_month = report_month_dt.month
    delta_base = report_month - 1
    monthly_series: dict[int, dict[str, dict[str, int | None]]] = {}

    for month in range(1, report_month + 1):
        delta = month - 1 - delta_base
        month_start = shift_month(start_dt, delta).strftime("%Y-%m-%d")
        month_end = shift_month(end_dt, delta).strftime("%Y-%m-%d")

        month_data: dict[str, dict[str, int | None]] = {}
        for lang in LANGUAGES:
            summary = fetch_summary_safe(
                property_ids[lang],
                month_start,
                month_end,
                context=f"monthly_series_{month}_{lang}",
            )
            is_missing = bool(summary.get("_missing"))
            month_data[lang] = {
                "users": None if is_missing else int(summary["users"]),
                "sessions": None if is_missing else int(summary["sessions"]),
                "pageviews": None if is_missing else int(summary["pageviews"]),
            }
        monthly_series[month] = month_data

    return monthly_series


def _normalize_text(text: str | None) -> str:
    """매핑 비교를 위한 문자열 정규화."""
    return " ".join((text or "").strip().split()).lower()


def _is_blank_cell(value) -> bool:
    """IA 셀 값이 비어있거나 공백/전각 공백인지 확인."""
    if value is None:
        return True
    text = str(value).strip()
    return not text or text == "　"


def _build_ia_title_maps() -> dict[str, dict[str, str]]:
    """
    IA 파일 기준으로 언어별 title->breadcrumb 매핑을 구성한다.
    A~C(depth)를 breadcrumb로 만들고, E~G(title) 값으로 키를 생성한다.
    """
    ia_path = TEMPLATES_DIR / IA_FILENAME
    if not ia_path.exists():
        return {lang: {} for lang in LANGUAGES}

    workbook = load_workbook(ia_path, data_only=True)
    sheet = workbook["IA_KR"] if "IA_KR" in workbook.sheetnames else workbook.worksheets[0]

    title_maps = {lang: {} for lang in LANGUAGES}
    current_depth = ["", "", ""]

    for row_idx in range(4, sheet.max_row + 1):
        row_values = [sheet.cell(row_idx, col_idx).value for col_idx in range(1, 8)]

        for depth_idx in range(3):
            raw_depth = row_values[depth_idx]
            if not _is_blank_cell(raw_depth):
                current_depth[depth_idx] = str(raw_depth).strip()
                for reset_idx in range(depth_idx + 1, 3):
                    if _is_blank_cell(row_values[reset_idx]):
                        current_depth[reset_idx] = ""

        breadcrumb_parts = [part for part in current_depth if part]
        if not breadcrumb_parts:
            continue
        breadcrumb = " > ".join(breadcrumb_parts)

        for lang in LANGUAGES:
            title_col = IA_LANG_COLS[lang] - 1
            title_raw = row_values[title_col]
            if _is_blank_cell(title_raw):
                continue
            key = _normalize_text(str(title_raw))
            if key:
                title_maps[lang][key] = breadcrumb

    return title_maps


def _simplify_for_match(text: str | None) -> str:
    """메인 페이지 동치 판별을 위한 공격적 정규화."""
    source = (text or "").lower()
    for token in (" ", "\t", "\n", "\r", "|", ">", "-", "_", "/", "\\"):
        source = source.replace(token, "")
    return source


def _is_main_equivalent(value: str | None) -> bool:
    """대상웰라이프/Daesang Wellife/大象Wellife 등 메인 페이지 동치 여부."""
    simplified = _simplify_for_match(value)
    if not simplified:
        return False
    return simplified in MAIN_EQUIVALENT_ALIASES


def _trim_title_suffix(title: str | None) -> str:
    """title 태그에서 사이트명 suffix를 제거해 라벨 오인식을 줄인다."""
    text = (title or "").strip()
    if not text:
        return ""
    return text.split("|")[0].strip()


def _fetch_slide6_top_pages(
    property_id: str,
    start_date: str,
    end_date: str,
    lang: str,
    ia_title_maps: dict[str, dict[str, str]],
    fetch_limit: int = 25,
    output_limit: int = 25,
) -> tuple[list[dict], dict[str, int]]:
    """
    6p용 top pages 데이터를 페이지 제목 및 화면 클래스(unifiedScreenClass) 기준으로 조회한다.
    - 조회 차원부터 제목 단일 기준으로 고정해 경로 영향 제거
    - 메인 동치는 제목 문자열 기준으로 합산
    """
    client = get_ga4_client()
    request = RunReportRequest(
        property=property_id,
        date_ranges=[DateRange(start_date=start_date, end_date=end_date)],
        dimensions=[Dimension(name="unifiedScreenClass")],
        metrics=[Metric(name="screenPageViews")],
        order_bys=[
            OrderBy(metric=OrderBy.MetricOrderBy(metric_name="screenPageViews"), desc=True)
        ],
        limit=fetch_limit,
    )
    response = client.run_report(request)

    aggregated: dict[str, dict[str, int | str]] = {}
    title_map = ia_title_maps.get(lang, {})
    stats = {
        "rows": 0,
        "main_merged_rows": 0,
        "not_set_rows": 0,
        "not_set_excluded_rows": 0,
        "ia_mapped_rows": 0,
        "ga4_text_fallback_rows": 0,
        "title_grouped_rows": 0,
    }

    for row in response.rows:
        page_title = row.dimension_values[0].value
        pageviews = int(row.metric_values[0].value)
        stats["rows"] += 1

        trimmed_title = _trim_title_suffix(page_title)
        title_key = _normalize_text(page_title)
        mapped_path = title_map.get(title_key)
        if not mapped_path and trimmed_title:
            mapped_path = title_map.get(_normalize_text(trimmed_title))

        if mapped_path:
            display_page = mapped_path
            stats["ia_mapped_rows"] += 1
        else:
            display_page = trimmed_title or page_title or "-"
            stats["ga4_text_fallback_rows"] += 1

        normalized_title = _normalize_text(trimmed_title or page_title)
        if normalized_title in {"(not set)", "not set"}:
            stats["not_set_rows"] += 1
            stats["not_set_excluded_rows"] += 1
            continue

        # 집계는 제목 기준(현재 로직 유지)
        title_rank_key = normalized_title
        main_candidate = trimmed_title or page_title
        is_main = _is_main_equivalent(main_candidate)
        if is_main:
            rank_key = "main_title"
            display_page = "메인"
            stats["main_merged_rows"] += 1
        else:
            rank_key = title_rank_key or "__untitled__"

        if rank_key in aggregated:
            aggregated[rank_key]["pageviews"] += pageviews
            # 같은 제목 그룹에서 IA 경로가 뒤늦게 발견되면 IA 경로로 승격
            if mapped_path and aggregated[rank_key]["page"] != "메인":
                aggregated[rank_key]["page"] = mapped_path
        else:
            aggregated[rank_key] = {"page": display_page, "pageviews": pageviews}

    stats["title_grouped_rows"] = len(aggregated)
    ranked = sorted(
        aggregated.values(),
        key=lambda item: int(item["pageviews"]),
        reverse=True,
    )
    result = [
        {"page": str(item["page"]), "pageviews": int(item["pageviews"])}
        for item in ranked[:output_limit]
    ]
    stats["aggregated_rows"] = len(ranked)
    return result, stats


def apply_slide6_top_pages_override(
    client_config: dict,
    start_date: str,
    end_date: str,
    data: dict,
) -> None:
    """6페이지 요구사항에 맞춰 언어별 top_pages를 재구성한다."""
    property_ids = client_config.get("ga4_property_ids", {})
    ia_title_maps = _build_ia_title_maps()
    for lang in LANGUAGES:
        property_id = property_ids.get(lang)
        if not property_id:
            continue
        top_pages, stats = _fetch_slide6_top_pages(
            property_id=property_id,
            start_date=start_date,
            end_date=end_date,
            lang=lang,
            ia_title_maps=ia_title_maps,
            fetch_limit=25,
            output_limit=25,
        )
        data[lang]["top_pages"] = top_pages
        write_log(
            "SLIDE6_TOP_PAGES "
            f"lang={lang} rows={stats['rows']} main_merged={stats['main_merged_rows']} "
            f"ia_mapped={stats['ia_mapped_rows']} "
            f"ga4_text_fallback={stats['ga4_text_fallback_rows']} "
            f"not_set={stats['not_set_rows']} not_set_excluded={stats['not_set_excluded_rows']} "
            f"title_grouped={stats['title_grouped_rows']} "
            f"aggregated={stats['aggregated_rows']} title_only_query=true"
        )


def write_log(message: str) -> None:
    """실행 로그를 logs/run_log.txt에 append 한다."""
    LOG_DIR.mkdir(parents=True, exist_ok=True)
    timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    with LOG_PATH.open("a", encoding="utf-8") as file:
        file.write(f"[{timestamp}] {message}\n")


def run_report(
    client_name: str,
    start_raw: str,
    end_raw: str,
    report_month_raw: str | None = None,
) -> dict:
    """전체 보고서 생성 흐름을 실행하고 결과 파일 경로를 반환한다."""
    start_dt = parse_date(start_raw)
    end_dt = parse_date(end_raw)
    if start_dt > end_dt:
        raise ValueError("시작일은 종료일보다 늦을 수 없습니다.")

    start_date = format_ga4_date(start_raw)
    end_date = format_ga4_date(end_raw)
    report_month_dt = parse_report_month(report_month_raw) if report_month_raw else end_dt

    client_config = load_client_config(client_name)
    data = collect_ga4_data(client_config, start_date, end_date)
    monthly_summary_series = collect_monthly_summary_series(
        client_config=client_config,
        start_dt=start_dt,
        end_dt=end_dt,
        report_month_dt=report_month_dt,
    )
    apply_slide6_top_pages_override(client_config, start_date, end_date, data)

    year = report_month_dt.year
    month = report_month_dt.month
    report_month = format_report_month(report_month_dt)
    client_output_dir = get_client_output_dir(client_name)
    client_output_dir.mkdir(parents=True, exist_ok=True)

    ppt_output_path = resolve_output_path(client_name, report_month, "pptx")
    if ppt_output_path.exists():
        raise FileExistsError(
            f"해당 월 산출물이 이미 존재합니다: {report_month} "
            f"(ppt={ppt_output_path.exists()})"
        )

    ppt_base_path, ppt_base_source = resolve_base_report_path(
        client_name=client_name,
        report_month_dt=report_month_dt,
        template_filename=client_config["ppt_template"],
        extension="pptx",
    )
    annual_baseline = load_annual_baseline(client_name, report_month_dt)
    runtime_ppt_base_path = apply_annual_baseline_to_ppt_base(
        ppt_base_path=ppt_base_path,
        client_output_dir=client_output_dir,
        report_month=report_month,
        annual_baseline=annual_baseline,
        monthly_summary_series=monthly_summary_series,
    )

    ppt_path = ppt_gen.write_report(
        template_path=runtime_ppt_base_path,
        output_dir=str(client_output_dir),
        client_name=client_name,
        year=year,
        month=month,
        data=data,
        start_date=start_date,
        end_date=end_date,
    )
    apply_growth_overrides_to_generated_ppt(
        ppt_path=ppt_path,
        report_month_dt=report_month_dt,
        monthly_summary_series=monthly_summary_series,
        annual_baseline=annual_baseline,
    )

    return {
        "ppt_path": ppt_path,
        "ppt_base_path": ppt_base_path,
        "ppt_base_source": ppt_base_source,
        "report_month": report_month,
        "annual_baseline_loaded": True,
    }


def main() -> None:
    """CLI 진입점."""
    args = parse_args()

    try:
        result = run_report(args.client, args.start, args.end, args.report_month)
        write_log(
            f"SUCCESS client={args.client} start={args.start} end={args.end} "
            f"report_month={result['report_month']} "
            f"ppt={result['ppt_path']} "
            f"ppt_base={result['ppt_base_source']}:{result['ppt_base_path']} "
            f"annual_baseline_loaded={result['annual_baseline_loaded']}"
        )
        print("보고서 생성이 완료되었습니다.")
        print(f"PPT: {result['ppt_path']}")
        print(f"PPT base: {result['ppt_base_path']} ({result['ppt_base_source']})")
    except Exception as exc:
        write_log(
            f"FAIL client={args.client} start={args.start} end={args.end} "
            f"report_month={args.report_month or '(from_end)'} error={exc}"
        )
        print(f"오류가 발생했습니다: {exc}")
        raise


if __name__ == "__main__":
    main()
